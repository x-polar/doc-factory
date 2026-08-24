"""kori_dark_measure.json → 네이티브 pptx (옵션 B 파일럿, 다크 1장)
모든 도형은 실제 pptx 오브젝트: 존=점선 라운드 사각형, 노드=라운드 사각형,
연결선=자유형(freeform), 텍스트=텍스트박스, 아이콘=투명 PNG.
"""
import json, re, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

d = json.load(open('/tmp/kori_dark_measure.json'))
BG = (0, 0, 0)
TEXT = (237, 237, 237)
MINT = (0x22, 0xFF, 0xBD)
AMBER = (0xFF, 0xB2, 0x24)
GRAY = (0x8A, 0x8A, 0x8A)

def _parse_color(s):
    """rgb()/rgba()/color(srgb r g b / a) → ((r,g,b) 0..255, alpha 0..1)"""
    if not s:
        return None, 0
    m = re.findall(r'[\d.]+', s)
    if not m:
        return None, 0
    vals = [float(v) for v in m]
    if s.strip().startswith('color('):
        # color(srgb 0..1 0..1 0..1 / a?)
        rgb = tuple(int(round(v * 255)) for v in vals[:3])
        a = vals[3] if len(vals) > 3 else 1.0
        return rgb, a
    rgb = tuple(int(v) for v in vals[:3])
    a = vals[3] if len(vals) > 3 else 1.0
    return rgb, a

def css_rgb(s, default=(255, 255, 255)):
    rgb, _ = _parse_color(s)
    return rgb or default

def css_rgba(s):
    return _parse_color(s)

def px2pt(px):
    return float(px.replace('px', '')) * 0.75

def set_alpha(fill_or_line_color_elem, alpha):
    """srgbClr 엘리먼트에 alpha 추가"""
    a = fill_or_line_color_elem
    for tag in ('a:alpha',):
        pass
    el = a.makeelement(qn('a:alpha'), {'val': str(int(alpha * 100000))})
    a.append(el)

def solid_fill_alpha(shape, rgb, alpha=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    if alpha < 1.0:
        spPr = shape.fill._xPr  # spPr element
        sf = spPr.find(qn('a:solidFill'))
        srgb = sf.find(qn('a:srgbClr')) if sf is not None else None
        if srgb is not None:
            set_alpha(srgb, alpha)

def line_alpha(shape, alpha):
    ln = shape.line._get_or_add_ln()
    sf = ln.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            set_alpha(srgb, alpha)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── 배경 ──
bgs = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bgs.fill.solid(); bgs.fill.fore_color.rgb = RGBColor(*BG)
bgs.line.fill.background()
bgs.shadow.inherit = False

def txbox(x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
          font='Pretendard', spacing=None, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = RGBColor(*color); r.font.name = font
        if spacing is not None:
            rPr = r._r.get_or_add_rPr(); rPr.set('spc', str(int(spacing * 100)))
    return tb

# ── 제목 ──
t = d['title']
txbox(t['x'], t['y'] - 0.06, t['w'], t['h'] + 0.15, t['text'],
      px2pt(t['size']), css_rgb(t['color']), bold=True, font='Tomorrow')

# ── 존 (점선 라운드 사각형) ──
for z in d['zones']:
    zs = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(z['x']), Inches(z['y']),
                                Inches(z['w']), Inches(z['h']))
    zs.adjustments[0] = 0.045
    zs.shadow.inherit = False
    rgb, a = css_rgba(z['bg'])
    if rgb and a > 0:
        solid_fill_alpha(zs, rgb, min(a, 1.0))
    else:
        zs.fill.background()
    brgb, ba = css_rgba(z['border'])
    zs.line.color.rgb = RGBColor(*(brgb or GRAY))
    zs.line.width = Pt(1.2)
    zs.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if ba < 1:
        line_alpha(zs, max(ba, 0.8))
    # 존 라벨 (경계 위 칩)
    if z['label'] and z['lblRect']:
        lr = z['lblRect']
        lbl = txbox(lr['x'], lr['y'], lr['w'] + 0.35, lr['h'], z['label'],
                    px2pt(z['lblSize']), css_rgb(z['lblColor']), bold=True,
                    font='Pretendard', spacing=0.45)
        # 배경색으로 경계선 위에 얹힘 (dg-zone-label은 배경을 슬라이드색으로 가짐)
        lbl.fill.solid(); lbl.fill.fore_color.rgb = RGBColor(*BG)
        lbl.line.fill.background()
        lbl.text_frame.word_wrap = False

# ── 자유 라벨 (앰버 밴드 배경 + 앰버 배지 + Tool Invocation / Grounding Context) ──
for L in d['labels']:
    rgb, a = css_rgba(L['bg'])
    text = re.sub(r'<[^>]+>', '', L['html']).strip()
    if not text:
        # 배경 밴드 (앰버): 반투명 채움 + 점선 앰버 테두리
        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(L['x']), Inches(L['y']),
                                      Inches(L['w']), Inches(L['h']))
        band.adjustments[0] = 0.03
        band.shadow.inherit = False
        solid_fill_alpha(band, AMBER, 0.055)
        band.line.color.rgb = RGBColor(*AMBER)
        band.line.width = Pt(1.0)
        band.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        line_alpha(band, 0.55)
        continue
    is_badge = a > 0.5  # 앰버 라벨
    align = {'center': PP_ALIGN.CENTER}.get(L['align'], PP_ALIGN.LEFT)
    text = ' '.join(text.split()) if is_badge else text
    ext = 0.10 if is_badge else 0.55
    tb = txbox(L['x'] - ext / 2, L['y'], L['w'] + ext, L['h'] + 0.02, text if is_badge else L['html'],
               px2pt(L['size']), css_rgb(L['color']),
               bold=(L['weight'] in ('700', 'bold')), align=align,
               font='Pretendard',
               spacing=0.5 if is_badge else None, line_spacing=0.95)
    tb.text_frame.word_wrap = False
    if is_badge:
        tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(*rgb)
        tb.line.fill.background()
        tb.text_frame.margin_left = tb.text_frame.margin_right = Pt(4)
        tb.text_frame.margin_top = tb.text_frame.margin_bottom = Pt(1)

# ── 연결선 (freeform: 시작→끝 폴리라인) ──
def add_edge(e):
    pts = e['pts']
    stroke = css_rgb(e['stroke'], GRAY)
    dashed = e['dash'] not in ('none', '', None)
    is_mint = 'e--accent' in (e['cls'] or '')
    color = MINT if is_mint else stroke
    for i in range(len(pts) - 1):
        (x1, y1), (x2, y2) = pts[i], pts[i + 1]
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = RGBColor(*color)
        conn.line.width = Pt(1.2)
        conn.shadow.inherit = False
        if dashed:
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if not is_mint:
            _, a = css_rgba(e['stroke'])
            if a and a < 1:
                line_alpha(conn, max(a, 0.85))
        # 화살촉 (마지막 세그먼트에만)
        if e['marker'] and i == len(pts) - 2:
            ln = conn.line._get_or_add_ln()
            he = ln.makeelement(qn('a:tailEnd'), {'type': 'arrow', 'w': 'med', 'len': 'med'})
            ln.append(he)

for e in d['edges']:
    add_edge(e)

# ── SVG 텍스트 라벨 (OpenAI Compatible API, Load) ──
for st in d['svgtexts']:
    size = px2pt(st['size'])
    w = len(st['text']) * size / 72 * 0.62 + 0.1
    h = size / 72 * 1.6
    x = st['x'] - (w if st['anchor'] == 'end' else w / 2 if st['anchor'] == 'middle' else 0)
    align = {'end': PP_ALIGN.RIGHT, 'middle': PP_ALIGN.CENTER}.get(st['anchor'], PP_ALIGN.LEFT)
    txbox(x, st['y'] - h * 0.75, w, h, st['text'], size, css_rgb(st['fill'], GRAY), align=align)

# ── 노드 ──
for n in d['nodes']:
    ns = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(n['x']), Inches(n['y']),
                                Inches(n['w']), Inches(n['h']))
    ns.adjustments[0] = min(0.5, 0.09 / min(n['w'], n['h']))
    ns.shadow.inherit = False
    if n['ghost']:
        ns.fill.background()
    else:
        # 그라데이션 근사: 어두운 차콜 단색
        solid_fill_alpha(ns, (23, 23, 23))
    brgb, ba = css_rgba(n['border'])
    ns.line.color.rgb = RGBColor(*(brgb or GRAY))
    ns.line.width = Pt(1.0)
    if n['ghost']:
        ns.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if ba < 1:
        line_alpha(ns, ba)
    # 아이콘
    if n.get('icoKey') and n['icoRect']:
        ir = n['icoRect']
        iy = ir['y']
        if n['title'] and n['tRect']:
            lines = n['title'].count('\n') + 1
            line_h = n['tRect']['h'] / lines
            iy = n['tRect']['y'] + (line_h - ir['h']) / 2
        slide.shapes.add_picture(f"/tmp/kori_icons/{n['icoKey']}.png",
                                 Inches(ir['x']), Inches(iy), Inches(ir['w']), Inches(ir['h']))
    # 제목
    if n['title'] and n['tRect']:
        tr = n['tRect']
        ttb = txbox(tr['x'], tr['y'], tr['w'] + 0.4, tr['h'] + 0.02, n['title'],
              px2pt(n['tSize']), css_rgb(n['tColor']), bold=(n['tWeight'] == '700'),
              font='Pretendard', line_spacing=0.98)
        ttb.text_frame.word_wrap = False
        ttb.text_frame.auto_size = None
    # 설명 (dg-s)
    if n['sub']:
        sy = (n['tRect']['y'] + n['tRect']['h'] + 0.03) if n['tRect'] else n['y'] + 0.3
        txbox(n['x'] + 0.12, sy, n['w'] - 0.2, 0.25, n['sub'],
              px2pt(n['sSize']), css_rgb(n['sColor'], GRAY))
    # 리스트 (dg-items) — 민트 체크 불릿
    if n['items']:
        iy = (n['tRect']['y'] + n['tRect']['h'] + 0.045) if n['tRect'] else n['y'] + 0.35
        isz = px2pt(n['itemSize'] or '9px')
        tb = slide.shapes.add_textbox(Inches(n['x'] + 0.13), Inches(iy),
                                      Inches(n['w'] - 0.24), Inches(n['h'] - (iy - n['y']) - 0.06))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, it in enumerate(n['items']):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.15
            r1 = p.add_run(); r1.text = '✓ '
            r1.font.size = Pt(isz); r1.font.color.rgb = RGBColor(*MINT); r1.font.name = 'Pretendard'
            r2 = p.add_run(); r2.text = it
            r2.font.size = Pt(isz); r2.font.color.rgb = RGBColor(0xC9, 0xC9, 0xC9)
            r2.font.name = 'Pretendard'

# ── 푸터 ──
f = d['footer']
if f.get('logo'):
    lg = f['logo']
    src = '/tmp/kori_icons/footer_logo.png'
    slide.shapes.add_picture(src, Inches(lg['x']), Inches(lg['y']), Inches(lg['w']), Inches(lg['h']))
if f.get('pg'):
    p = f['pg']
    txbox(p['x'] - 0.1, p['y'] - 0.02, p['w'] + 0.2, p['h'] + 0.05, p['text'],
          px2pt(p['size']), css_rgb(p['color'], GRAY), font='Tomorrow', align=PP_ALIGN.RIGHT)

out = 'docs/20260822_kori_architecture/output/20260822_kori_architecture_main-dark_native_v1.pptx'
prs.save(out)
print('saved:', out)
