#!/usr/bin/env python3
"""storyboard(슬라이드별 .md) + 브랜드 테마 -> .pptx 생성.

문서의 brief.md frontmatter가 브랜드를 고르고(brands/<brand>/theme.yaml), 필요하면
theme: 블록으로 토큰을 오버라이드한다. 마스터 pptx 없이 토큰으로 프로그래밍
스타일링한다(마스터가 생기면 그 위에 얹도록 확장). layout 정의는
reference/layout-catalog.md, 브랜드 구조는 brands/README.md 참고.

사용법:
  python src/lib/deckbuilder.py <문서폴더> [-o out.pptx] [--brand NAME] [--theme path]
예:
  python src/lib/deckbuilder.py docs/20260629_proposal
"""
import argparse
import csv
import glob
import os
import re

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))


# ----- 입력 파싱 ---------------------------------------------------------------

def deep_merge(base, over):
    """over의 값으로 base를 재귀 병합(dict는 깊게, 그 외는 덮어씀)."""
    out = dict(base)
    for k, v in (over or {}).items():
        if isinstance(v, dict) and isinstance(out.get(k), dict):
            out[k] = deep_merge(out[k], v)
        else:
            out[k] = v
    return out


def load_doc_meta(doc_dir):
    """문서 메타(title/brand/version/theme)를 brief.md frontmatter에서 읽는다."""
    path = os.path.join(doc_dir, "brief.md")
    if not os.path.exists(path):
        return {}
    text = open(path, encoding="utf-8").read()
    m = re.match(r"^---\n(.*?)\n---\n", text, re.S)
    return (yaml.safe_load(m.group(1)) or {}) if m else {}


def resolve_brand_dir(brand):
    brand_dir = os.path.join(ROOT, "brands", brand or "_default")
    if not os.path.isdir(brand_dir):
        brand_dir = os.path.join(ROOT, "brands", "_default")
    return brand_dir


def resolve_theme(doc_dir, brand_override=None, theme_path=None):
    """brief.md frontmatter -> 브랜드 theme 로드 -> theme: 오버라이드 병합.
    반환: (theme dict, brand_dir, meta dict)."""
    meta = load_doc_meta(doc_dir)
    brand = brand_override or meta.get("brand") or "_default"
    brand_dir = resolve_brand_dir(brand)
    theme_file = theme_path or os.path.join(brand_dir, "theme.yaml")
    theme = yaml.safe_load(open(theme_file, encoding="utf-8"))
    theme = deep_merge(theme, meta.get("theme"))
    return theme, brand_dir, meta


def load_slides(storyboard_dir):
    """storyboard/*.md 를 파일명 순으로 읽어 슬라이드 dict 목록 반환."""
    slides = []
    for path in sorted(glob.glob(os.path.join(storyboard_dir, "*.md"))):
        text = open(path, encoding="utf-8").read()
        m = re.match(r"^---\n(.*?)\n---\n(.*)$", text, re.S)
        meta = (yaml.safe_load(m.group(1)) if m else {}) or {}
        rest = m.group(2) if m else text
        parts = re.split(r"^##\s*notes\s*$", rest, maxsplit=1, flags=re.M)
        lines = parts[0].splitlines()
        body = [l.strip()[1:].strip() for l in lines if l.strip().startswith("-")]
        free = [l.strip() for l in lines
                if l.strip() and not l.strip().startswith("-")]
        notes = parts[1].strip() if len(parts) > 1 else ""
        slides.append({**meta, "body": body, "free": free,
                       "notes": notes, "_file": os.path.basename(path)})
    return slides


def read_csv_chart(path):
    """첫 행 헤더·첫 열 카테고리 CSV -> (categories, [(series_name, values)])."""
    with open(path, encoding="utf-8") as f:
        rows = [r for r in csv.reader(f) if r]
    header = rows[0]
    series_names = header[1:]
    categories = [r[0] for r in rows[1:]]
    series = []
    for i, name in enumerate(series_names, start=1):
        series.append((name, [float(r[i]) for r in rows[1:]]))
    return categories, series


# ----- 렌더 헬퍼 ---------------------------------------------------------------

class Deck:
    def __init__(self, theme):
        self.t = theme
        self.c = theme["colors"]
        self.f = theme["fonts"]
        self.s = theme["sizes"]
        self.sl = theme["slide"]
        self.logo = theme.get("logo", {})   # {cover: path, footer: path} (브랜드 상대경로)
        self.brand_dir = None               # build()에서 자산 해석 기준으로 설정
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.sl["width"])
        self.prs.slide_height = Inches(self.sl["height"])
        self.blank = self.prs.slide_layouts[6]  # 기본 템플릿의 Blank 레이아웃

    def rgb(self, key):
        return RGBColor.from_string(self.c[key])

    def col(self, key, default_key=None):
        """토큰 색을 가져오되 없으면 default_key 토큰으로 폴백(하위호환)."""
        if key in self.c:
            return self.c[key]
        return self.c[default_key]

    # ----- 색/폰트 유틸 -----
    @staticmethod
    def _has_cjk(s):
        return any('가' <= ch <= '힣' or '㄰' <= ch <= '㆏'
                   or '一' <= ch <= '鿿' for ch in (s or ""))

    def _heading_font(self, text):
        """제목에 한글/한자가 있으면 CJK 지원 폰트로, 아니면 디스플레이 heading으로.
        제목 단위로 폰트를 통일해 혼용 시 들쭉날쭉함을 줄인다."""
        if self._has_cjk(text):
            return self.f.get("heading_cjk", self.f["body"])
        return self.f["heading"]

    @staticmethod
    def _shade(hexstr, toward_white=0.0, toward_black=0.0):
        r, g, b = int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16)
        if toward_white:
            r = int(r + (255 - r) * toward_white); g = int(g + (255 - g) * toward_white); b = int(b + (255 - b) * toward_white)
        if toward_black:
            r = int(r * (1 - toward_black)); g = int(g * (1 - toward_black)); b = int(b * (1 - toward_black))
        return "%02X%02X%02X" % (r, g, b)

    @staticmethod
    def _is_dark(hexstr):
        r, g, b = int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) < 128

    def _bg_hex(self):
        return self.col("background", None) if "background" in self.c else "FFFFFF"

    def _card_bg(self):
        if "card_bg" in self.c:
            return self.c["card_bg"]
        bg = self._bg_hex()
        return self._shade(bg, toward_white=0.10) if self._is_dark(bg) else self._shade(bg, toward_black=0.05)

    def _accent_rule(self, slide, y, width=1.6):
        """제목 아래 액센트 가로 룰(시각적 구조)."""
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(self.sl["margin"]),
                                     Inches(y), Inches(width), Pt(3))
        shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor.from_string(self.c["accent"])
        shp.line.fill.background(); shp.shadow.inherit = False
        return shp

    def _add_text(self, slide, text, left, top, width, height, *,
                  size, color, bold=False, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, font=None):
        box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font or self.f["body"]
        run.font.color.rgb = RGBColor.from_string(color)
        return box

    def _bg(self, slide, hexstr):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(hexstr)

    def _notes(self, slide, text):
        if text:
            slide.notes_slide.notes_text_frame.text = text

    def _add_logo(self, slide, kind):
        """theme.logo[kind] 로고를 배치. cover=상단 크게, footer=하단 작게."""
        rel = self.logo.get(kind)
        if not rel or not self.brand_dir:
            return
        path = os.path.join(self.brand_dir, rel)
        if not os.path.exists(path):
            return
        m = self.sl["margin"]
        if kind == "cover":
            slide.shapes.add_picture(path, Inches(m), Inches(m), width=Inches(3.4))
        else:  # footer
            slide.shapes.add_picture(path, Inches(m),
                                     Inches(self.sl["height"] - 0.55), width=Inches(1.5))

    # --- layout 별 렌더 ---
    def render(self, sp):
        layout = sp.get("layout", "title+body")
        fn = getattr(self, "_l_" + layout.replace("+", "_").replace("-", "_"), None)
        if fn is None:
            fn = self._l_title_body  # 알 수 없는 layout은 본문형으로 폴백
        slide = self.prs.slides.add_slide(self.blank)
        fn(slide, sp)
        # 로고: 표지는 상단 크게, 나머지는 푸터에 작게
        self._add_logo(slide, "cover" if layout == "title" else "footer")
        self._notes(slide, sp.get("notes", ""))
        return slide

    def _content_bg(self, slide):
        # 콘텐츠 슬라이드 배경(라이트=흰색, 다크 브랜드=어두운 색). 기본 흰색.
        self._bg(slide, self.col("background", None) if "background" in self.c else "FFFFFF")

    def _title_band(self, slide, text):
        """콘텐츠 제목 + 액센트 룰. 본문 시작 y(인치)를 반환(2줄 제목이면 더 내림)."""
        m = self.sl["margin"]
        self._add_text(slide, text, m, m - 0.05, self.sl["width"] - 2 * m, 1.4,
                       size=self.s["title"], color=self.col("title_ink", "primary"),
                       bold=True, font=self._heading_font(text))
        # CJK는 폭 2로 환산해 줄바꿈 추정 → 2줄이면 룰/본문을 아래로
        units = sum(2 if self._has_cjk(ch) else 1 for ch in text)
        two_line = units > 56
        rule_y = 1.95 if two_line else 1.55
        self._accent_rule(slide, rule_y)
        return rule_y + 0.3

    def _bullets(self, slide, items, top, *, left=None, width=None, size=None):
        """액센트 불릿 목록을 그린다."""
        m = self.sl["margin"]
        left = self.sl["margin"] if left is None else left
        width = (self.sl["width"] - 2 * m) if width is None else width
        size = self.s["body"] if size is None else size
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                       Inches(self.sl["height"] - top - 0.7))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            b = para.add_run(); b.text = "▪  "
            b.font.size = Pt(size); b.font.name = self.f["body"]
            b.font.color.rgb = RGBColor.from_string(self.c["accent"])
            r = para.add_run(); r.text = line
            r.font.size = Pt(size); r.font.name = self.f["body"]
            r.font.color.rgb = RGBColor.from_string(self.c["text"])
            para.space_after = Pt(10)
        return box

    def _l_title(self, slide, sp):
        self._bg(slide, self.col("bg_cover", "primary"))
        w, h = self.sl["width"], self.sl["height"]
        on = self.col("on_cover", "on_primary")
        self._add_text(slide, sp["title"], 1.0, h / 2 - 1.0, w - 2.0, 1.5,
                       size=self.s["title"] + 8, color=on,
                       bold=True, align=PP_ALIGN.LEFT, font=self._heading_font(sp["title"]))
        self._accent_rule(slide, h / 2 + 0.15)
        sub = " · ".join(sp.get("free", []))
        if sub:
            self._add_text(slide, sub, 1.0, h / 2 + 0.35, w - 2.0, 1.0,
                           size=self.s["subtitle"], color=on,
                           font=self.f["body"])

    def _l_section(self, slide, sp):
        self._bg(slide, self.col("bg_section", "accent"))
        w, h = self.sl["width"], self.sl["height"]
        self._add_text(slide, sp["title"], 1.0, h / 2 - 0.6, w - 2.0, 1.2,
                       size=self.s["title"] + 4, color=self.col("on_section", "on_primary"),
                       bold=True, anchor=MSO_ANCHOR.MIDDLE, font=self._heading_font(sp["title"]))

    def _l_title_body(self, slide, sp):
        self._content_bg(slide)
        top = self._title_band(slide, sp["title"])
        self._bullets(slide, sp.get("body", []), top)

    def _l_columns(self, slide, sp):
        """제목 + N개 카드(columns: [{heading, items:[...]}, ...])."""
        self._content_bg(slide)
        top = self._title_band(slide, sp["title"])
        cols = sp.get("columns", []) or []
        n = max(1, len(cols))
        m, gap = self.sl["margin"], 0.3
        total_w = self.sl["width"] - 2 * m
        cw = (total_w - gap * (n - 1)) / n
        ch = self.sl["height"] - top - 0.7
        card_bg = self._card_bg()
        for i, c in enumerate(cols):
            x = m + i * (cw + gap)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(x), Inches(top), Inches(cw), Inches(ch))
            card.fill.solid(); card.fill.fore_color.rgb = RGBColor.from_string(card_bg)
            card.line.color.rgb = RGBColor.from_string(self.c["accent"]); card.line.width = Pt(0.75)
            card.shadow.inherit = False
            tf = card.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.2)
            head = c.get("heading", "")
            p0 = tf.paragraphs[0]
            rh = p0.add_run(); rh.text = head; rh.font.bold = True
            rh.font.size = Pt(self.s["body"] + 3); rh.font.name = self._heading_font(head)
            rh.font.color.rgb = RGBColor.from_string(self.c["accent"])
            for item in c.get("items", []) or []:
                p = tf.add_paragraph(); p.space_before = Pt(6)
                rb = p.add_run(); rb.text = "▪ "
                rb.font.size = Pt(self.s["caption"] + 1); rb.font.name = self.f["body"]
                rb.font.color.rgb = RGBColor.from_string(self.c["accent"])
                rt = p.add_run(); rt.text = item
                rt.font.size = Pt(self.s["caption"] + 1); rt.font.name = self.f["body"]
                rt.font.color.rgb = RGBColor.from_string(self.c["text"])

    def _cell(self, cell, text, *, bold=False, color="FFFFFF", fill=None):
        if fill:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(fill)
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = text
        r.font.size = Pt(self.s["caption"]); r.font.bold = bold
        r.font.name = self.f["body"]; r.font.color.rgb = RGBColor.from_string(color)

    def _l_table(self, slide, sp):
        """제목 + 표(table: {headers:[...], rows:[[...],...]})."""
        self._content_bg(slide)
        top = self._title_band(slide, sp["title"])
        t = sp.get("table", {}) or {}
        headers = t.get("headers", []); rows = t.get("rows", [])
        ncols = max(1, len(headers)); nrows = len(rows) + 1
        m = self.sl["margin"]; w = self.sl["width"] - 2 * m
        h = min(self.sl["height"] - top - 0.6, 0.4 * nrows + 0.2)
        gtbl = slide.shapes.add_table(nrows, ncols, Inches(m), Inches(top),
                                      Inches(w), Inches(h)).table
        gtbl.first_row = False  # python-pptx 기본 헤더 스타일 끄고 직접 칠함
        # 액센트(헤더 배경)가 밝으면 어두운 글자, 어두우면 흰 글자로 대비 확보
        header_ink = self.col("bg_cover", "primary") if not self._is_dark(self.c["accent"]) else self.c["on_primary"]
        for j, htext in enumerate(headers):
            self._cell(gtbl.cell(0, j), str(htext), bold=True, color=header_ink, fill=self.c["accent"])
        body_fill = self._card_bg()
        for i, row in enumerate(rows, start=1):
            for j in range(ncols):
                val = row[j] if j < len(row) else ""
                self._cell(gtbl.cell(i, j), str(val), color=self.c["text"], fill=body_fill)

    def _l_stat(self, slide, sp):
        """큰 숫자 강조(stats: [{value, label}, ...]) + 선택 본문."""
        self._content_bg(slide)
        top = self._title_band(slide, sp["title"]) if sp.get("title") else self.sl["margin"]
        stats = sp.get("stats", []) or []
        n = max(1, len(stats)); m = self.sl["margin"]
        total = self.sl["width"] - 2 * m; cw = total / n
        y = top + 0.5
        for i, s in enumerate(stats):
            x = m + i * cw
            self._add_text(slide, str(s.get("value", "")), x, y, cw, 1.5,
                           size=54, color=self.c["accent"], bold=True,
                           align=PP_ALIGN.CENTER, font=self.f["heading"])
            self._add_text(slide, s.get("label", ""), x, y + 1.5, cw, 0.9,
                           size=self.s["body"], color=self.c["text"],
                           align=PP_ALIGN.CENTER)
        if sp.get("body"):
            self._bullets(slide, sp["body"], y + 2.7)

    def _l_two_col(self, slide, sp):
        """제목 + 2단(left:[...], right:[...] 또는 right_image)."""
        self._content_bg(slide)
        top = self._title_band(slide, sp["title"])
        m, gap = self.sl["margin"], 0.4
        w = (self.sl["width"] - 2 * m - gap) / 2
        self._bullets(slide, sp.get("left", []), top, left=m, width=w)
        rimg = sp.get("_right_image_abspath")
        if rimg and os.path.exists(rimg):
            slide.shapes.add_picture(rimg, Inches(m + w + gap), Inches(top), width=Inches(w))
        else:
            self._bullets(slide, sp.get("right", []), top, left=m + w + gap, width=w)

    def _l_title_chart(self, slide, sp):
        self._content_bg(slide)
        self._title_band(slide, sp["title"])
        chart_path = sp.get("_chart_abspath")
        if not chart_path or not os.path.exists(chart_path):
            self._add_text(slide, f"[차트 누락: {sp.get('chart')}]",
                           self.sl["margin"], 1.8, 6, 0.5,
                           size=self.s["caption"], color=self.c["muted"])
            return
        categories, series = read_csv_chart(chart_path)
        data = CategoryChartData()
        data.categories = categories
        for name, vals in series:
            data.add_series(name, vals)
        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(self.sl["margin"]), Inches(1.8),
            Inches(self.sl["width"] - 2 * self.sl["margin"]), Inches(4.6), data)
        chart = gframe.chart
        chart.has_title = False
        # 다크 배경에서도 축·범례가 보이도록 차트 글자색을 text 토큰으로
        chart.font.color.rgb = RGBColor.from_string(self.c["text"])
        for i, plot_series in enumerate(chart.series):
            palette = self.c["series"]
            plot_series.format.fill.solid()
            plot_series.format.fill.fore_color.rgb = \
                RGBColor.from_string(palette[i % len(palette)])

    def _l_title_image(self, slide, sp):
        self._content_bg(slide)
        self._title_band(slide, sp["title"])
        img = sp.get("_image_abspath")
        if img and os.path.exists(img):
            slide.shapes.add_picture(img, Inches(self.sl["margin"]), Inches(1.8),
                                     height=Inches(4.6))
        else:
            self._add_text(slide, f"[이미지 누락: {sp.get('image')}]",
                           self.sl["margin"], 1.8, 6, 0.5,
                           size=self.s["caption"], color=self.c["muted"])

    def save(self, path):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)


# ----- 엔트리포인트 ------------------------------------------------------------

def build(doc_dir, out=None, brand=None, theme_path=None):
    theme, brand_dir, meta = resolve_theme(doc_dir, brand, theme_path)
    slides = load_slides(os.path.join(doc_dir, "storyboard"))
    # chart/image 경로를 문서폴더 기준 절대경로로 해석
    for sp in slides:
        if sp.get("chart"):
            sp["_chart_abspath"] = os.path.join(doc_dir, sp["chart"])
        if sp.get("image"):
            sp["_image_abspath"] = os.path.join(doc_dir, sp["image"])
        if sp.get("right_image"):
            sp["_right_image_abspath"] = os.path.join(doc_dir, sp["right_image"])
    deck = Deck(theme)
    deck.brand_dir = brand_dir  # 자산(로고 등) 해석 기준
    for sp in slides:
        deck.render(sp)
    if out is None:
        name = os.path.basename(os.path.normpath(doc_dir))
        version = meta.get("version")
        fname = f"{name}_{version}.pptx" if version else f"{name}.pptx"
        out = os.path.join(doc_dir, "output", fname)
    deck.save(out)
    print(f"생성: {out}  (슬라이드 {len(slides)}장, 브랜드: {os.path.basename(brand_dir)})")
    return out


def main():
    ap = argparse.ArgumentParser(description="storyboard -> pptx")
    ap.add_argument("doc_dir", help="문서 폴더 (storyboard/ 포함)")
    ap.add_argument("-o", "--out", help="출력 pptx 경로")
    ap.add_argument("--brand", help="brands/<NAME> 강제 지정(brief.md frontmatter보다 우선)")
    ap.add_argument("--theme", help="theme.yaml 경로 직접 지정")
    args = ap.parse_args()
    build(args.doc_dir, args.out, args.brand, args.theme)


if __name__ == "__main__":
    main()
