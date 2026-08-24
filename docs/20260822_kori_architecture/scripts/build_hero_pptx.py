#!/usr/bin/env python3
"""KORI Answers hero diagram → native-shape .pptx (편집용 재작도, 방식 B).

hero.html(v24, 커밋 2ad4e438 기준)의 CSS 2.5D 좌표를 python-pptx freeform/
도형으로 재구성한다. 모든 요소가 개별(그룹 없는) 네이티브 도형이라 PowerPoint
/Keynote에서 자유롭게 드래그·수정 가능.

근사 사항 (CSS → pptx 한계):
- iso matrix(0.866,.5,-.866,.5)는 freeform 정점 좌표에 직접 베이크 (완전 동일)
- 존 태그의 iso 눕힘은 rotation 30°로 근사 (pptx는 skew 미지원)
- 발광(glow)·color-mix 그라데이션은 단색 근사
- 아이콘 글리프는 단순 도형 근사 또는 생략

렌더 좌표계: dg-canvas 원점 = 슬라이드 (0.695in, 1.352in) [PNG 실측]
SVG edges viewBox 1190x540 → 캔버스 11.933 x 5.448in
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import math, os

# ── 캔버스 → 슬라이드 오프셋 ─────────────────────────────
OX, OY = 0.695, 1.352
CANVAS_W, CANVAS_H = 13.333 - 1.4, 6.8 - 1.352
SVGX, SVGY = CANVAS_W / 1190.0, CANVAS_H / 540.0

# ── 브랜드 토큰 ──────────────────────────────────────────
BG      = (0x00, 0x00, 0x00)
TEXT    = (0xED, 0xED, 0xED)
ACCENT  = (0x22, 0xFF, 0xBD)
MUTED   = (0x8A, 0x8A, 0x8A)
LAYER   = (0x2E, 0x9E, 0x86)
STORE   = (0x6F, 0xE9, 0xCE)
MODEL   = (0xD5, 0xD5, 0xD5)
XCUT    = (0xFF, 0xB2, 0x24)
DUCT    = (0x5A, 0x6B, 0x66)
DICT    = (0xB5, 0xA4, 0x6A)
MUTEBOX = (0x4A, 0x4F, 0x4D)

def mix(c, pct, base=BG):
    """color-mix(in srgb, c pct%, base) 근사."""
    return RGBColor(*(round(a * pct + b * (1 - pct)) for a, b in zip(c, base)))

def rgb(c):
    return RGBColor(*c)

# ── iso 수학 ─────────────────────────────────────────────
def iso(u, v):
    """matrix(0.866,0.5,-0.866,0.5): 로컬 (u,v) → 화면 delta."""
    return 0.866 * (u - v), 0.5 * (u + v)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

shapes = slide.shapes

# ── 배경 ────────────────────────────────────────────────
bgfill = slide.background.fill
bgfill.solid()
bgfill.fore_color.rgb = rgb(BG)

def freeform(pts, fill=None, line=None, line_w=1.0, dash=None, close=True):
    """pts: [(x,y) in slide inches]. dash: 'dash' | 'sysDot' 등."""
    fb = shapes.build_freeform(Emu(int(pts[0][0] * 914400)), Emu(int(pts[0][1] * 914400)), scale=1)
    fb.add_line_segments([(Emu(int(x * 914400)), Emu(int(y * 914400))) for x, y in pts[1:]], close=close)
    sh = fb.convert_to_shape()
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
        if dash:
            ln = sh.line._get_or_add_ln()
            d = ln.makeelement(qn('a:prstDash'), {'val': dash})
            ln.append(d)
    sh.shadow.inherit = False
    return sh

def plane(left, top, w, h, tint, fill_pct, border_pct, z_name):
    """dg-iso-plane: 대시 보더 평행사변형."""
    ax, ay = OX + left, OY + top
    pts = [(ax + iso(u, v)[0], ay + iso(u, v)[1]) for u, v in [(0, 0), (w, 0), (w, h), (0, h)]]
    sh = freeform(pts, fill=mix(tint, fill_pct), line=mix(tint, border_pct), line_w=1.2, dash='dash')
    sh.name = z_name
    return sh

def iso_box(left, top, sx, sy, hh, c, name, hero=False):
    """dg-iso-box: f-top/f-left/f-right 3면. 앵커 = 윗면 위 꼭짓점."""
    ax, ay = OX + left, OY + top
    # f-top
    top_pts = [(ax + iso(u, v)[0], ay + iso(u, v)[1]) for u, v in [(0, 0), (sx, 0), (sx, sy), (0, sy)]]
    # f-left: translate(-.866sy,.5sy) then (x,y)->(.866x,.5x+y); rect sx x hh
    tlx, tly = ax - 0.866 * sy, ay + 0.5 * sy
    left_pts = [(tlx + 0.866 * x, tly + 0.5 * x + y) for x, y in [(0, 0), (sx, 0), (sx, hh), (0, hh)]]
    # f-right: translate(.866(sx-sy),.5(sx+sy)) then (x,y)->(.866x,-.5x+y); rect sy x hh
    trx, try_ = ax + 0.866 * (sx - sy), ay + 0.5 * (sx + sy)
    right_pts = [(trx + 0.866 * x, try_ - 0.5 * x + y) for x, y in [(0, 0), (sy, 0), (sy, hh), (0, hh)]]
    border = mix(c, 0.85) if hero else mix(c, 0.60, base=(0x1F, 0x1F, 0x1F))
    lw = 1.4 if hero else 1.0
    s1 = freeform(left_pts, fill=mix(c, 0.52), line=border, line_w=lw);  s1.name = f"{name} · left"
    s2 = freeform(right_pts, fill=mix(c, 0.74, base=(0x0A, 0x0A, 0x0A)), line=border, line_w=lw); s2.name = f"{name} · right"
    s3 = freeform(top_pts, fill=mix(c, 0.28), line=border, line_w=lw);   s3.name = f"{name} · top"
    return s3

def iso_cyl(left, top, d, hh, c, name):
    """dg-iso-cyl → MSO CAN 도형 근사. 앵커 = 바운딩 좌상단."""
    sh = shapes.add_shape(MSO_SHAPE.CAN, Inches(OX + left), Inches(OY + top),
                          Inches(d), Inches(hh + d / 2))
    sh.fill.solid(); sh.fill.fore_color.rgb = mix(c, 0.55)
    sh.line.color.rgb = mix(c, 0.60, base=(0x1F, 0x1F, 0x1F)); sh.line.width = Pt(1.0)
    sh.shadow.inherit = False
    sh.name = name
    return sh

def edge(svg_pts, dashed=False, name="edge"):
    """SVG viewBox 좌표 폴리라인 → 캔버스 → 슬라이드."""
    pts = [(OX + x * SVGX, OY + y * SVGY) for x, y in svg_pts]
    sh = freeform(pts, fill=None, line=rgb(ACCENT), line_w=1.2,
                  dash='sysDash' if dashed else None, close=False)
    sh.name = name
    return sh

def textbox(left, top, text, size, color, bold=True, align=PP_ALIGN.LEFT,
            font="Pretendard", w=2.0, h=0.3, rot=0, name=None, spacing=0.0):
    tb = shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = color
    if spacing:
        r.font._rPr.set('spc', str(int(spacing * 100)))
    if rot:
        tb.rotation = rot
    tb.name = name or text
    return tb

def zonetag(left, top, text, tint, name):
    """iso 눕힘 라벨 → rot 30° 필 박스 근사. 앵커 = (translate -100%,-100% 무시하고
    대략 우하단 근사 배치: 폭 추정 후 좌상단 역산)."""
    w = 0.115 * len(text) * 0.75 / 0.72 * 0.72  # 7.5pt 대문자 근사 폭
    w = max(1.1, 0.13 * len(text) + 0.3)
    h = 0.17
    # translate(-100%,-100%): 회전 전 로컬 우하단이 앵커. rot 30° 근사 배치.
    ax, ay = OX + left, OY + top
    # 회전 중심(도형 중심)이 앵커에서 (-w/2,-h/2)를 30° 회전한 위치
    cx = ax + (-w / 2) * math.cos(math.radians(30)) - (-h / 2) * math.sin(math.radians(30))
    cy = ay + (-w / 2) * math.sin(math.radians(30)) + (-h / 2) * math.cos(math.radians(30))
    sh = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - w / 2), Inches(cy - h / 2),
                          Inches(w), Inches(h))
    sh.rotation = 30
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(tint)
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(7.5); r.font.bold = True; r.font.name = "Pretendard"
    r.font.color.rgb = rgb(BG)
    sh.name = name
    return sh

def leader(left, top, length, name):
    """수직 리더 라인 + 하단 도트."""
    ln = shapes.add_connector(1, Inches(OX + left), Inches(OY + top - length),
                              Inches(OX + left), Inches(OY + top))
    ln.line.color.rgb = mix(TEXT, 0.45); ln.line.width = Pt(1.0)
    ln.shadow.inherit = False; ln.name = f"{name} · leader"
    dot = shapes.add_shape(MSO_SHAPE.OVAL, Inches(OX + left - 0.028), Inches(OY + top - 0.028),
                           Inches(0.056), Inches(0.056))
    dot.fill.solid(); dot.fill.fore_color.rgb = mix(TEXT, 0.7)
    dot.line.fill.background(); dot.shadow.inherit = False
    dot.name = f"{name} · dot"

def callout(left, top, text, accent=False, anchor="bl", name=None):
    """dg-iso-callout: 세로 보더 + 텍스트. anchor='bl' 좌하단 / 'br' 우하단 /
    'tc' 상단중앙(translate(-50%,0))."""
    w = 0.085 * len(text) + 0.1
    h = 0.2
    ax, ay = OX + left, OY + top
    if anchor == "bl":
        bx, by = ax, ay - h
        line_x = ax
        tx = ax + 0.07
    elif anchor == "br":
        bx, by = ax - w, ay - h
        line_x = ax
        tx = ax - w
    else:  # tc
        bx, by = ax - w / 2, ay
        line_x = None
        tx = ax - w / 2 + 0.07
    if line_x is not None:
        vln = shapes.add_connector(1, Inches(line_x), Inches(by), Inches(line_x), Inches(by + h))
        vln.line.color.rgb = mix(TEXT, 0.45); vln.line.width = Pt(1.0)
        vln.shadow.inherit = False; vln.name = f"{name or text} · bar"
    col = rgb(ACCENT) if accent else rgb(TEXT)
    textbox(tx, by + 0.015, text, 8.5, col, bold=True, w=w, h=h, name=name or text)

# ═════════════════════════════════════════════════════════
# 1) 타이틀
textbox(0.7, 0.62, "KORI Answers", 20, rgb((0xFF, 0xFF, 0xFF)), font="Tomorrow",
        w=5, h=0.5, name="Title")

# 2) 배경 평면 (z0)
plane(5.15, 0.98, 2.1, 1.5, MODEL, 0.10, 0.55, "Plane · Inference Serving")
plane(7.55, 0.98, 1.7, 1.7, LAYER, 0.10, 0.55, "Plane · Tool Layer")
plane(9.55, 2.14, 2.0, 1.9, STORE, 0.10, 0.55, "Plane · Data Layer")
# X 앰버 밴드 (z1)
plane(6.6, 0.68, 2.7, 5.3, XCUT, 0.16, 0.70, "Plane · Deterministic Security")

# 3) 존 태그
zonetag(5.67, 2.78, "INFERENCE SERVING", MUTED, "Tag · INFERENCE SERVING")
zonetag(7.55, 2.68, "TOOL LAYER", LAYER, "Tag · TOOL LAYER")
zonetag(9.64, 4.09, "DATA LAYER", LAYER, "Tag · DATA LAYER")
zonetag(4.35, 4.60, "DETERMINISTIC SECURITY", XCUT, "Tag · DETERMINISTIC SECURITY")

# 4) 연결선 (SVG vb 1190x540)
edge([(92, 290), (126, 270), (196, 311)], name="Edge · User→Gateway")
edge([(296, 311), (368, 303)], name="Edge · Gateway→Agent")
edge([(441, 268), (441, 232), (540, 205)], name="Edge · Agent↕ModelGW")
edge([(497, 303), (662, 236), (729, 258)], name="Edge · Agent→Tools")
edge([(828, 225), (855, 241), (880, 262)], name="Edge · Tools→Stores")
edge([(786, 348), (800, 327)], dashed=True, name="Edge · Ingestion⇢Data")

# 5) End User — 사람 아이콘 근사 (머리 원 + 어깨 아치)
ux, uy = OX + 0.55, OY + 2.76
head = shapes.add_shape(MSO_SHAPE.OVAL, Inches(ux + 0.115), Inches(uy), Inches(0.17), Inches(0.17))
head.fill.background(); head.line.color.rgb = rgb(ACCENT); head.line.width = Pt(1.8)
head.shadow.inherit = False; head.name = "End User · head"
arc_pts = []
for i in range(13):
    th = math.pi * i / 12  # 180°..0° 상반원
    arc_pts.append((ux + 0.2 - 0.2 * math.cos(th), uy + 0.62 - 0.38 * math.sin(th)))
body = freeform(arc_pts, fill=None, line=rgb(ACCENT), line_w=1.8, close=False)
body.name = "End User · body"
textbox(ux + 0.2 - 0.4, uy + 0.67, "End User", 8, rgb(TEXT), align=PP_ALIGN.CENTER,
        w=0.8, h=0.2, name="End User · label")

# 6) G1 Compliance Gateway (muted box, mid h=0.26)
iso_box(2.55, 2.93, 0.6, 0.6, 0.26, MUTEBOX, "Compliance Gateway")
leader(2.55, 2.93, 0.85, "Compliance Gateway")
callout(2.55, 2.08, "Compliance Gateway", anchor="bl", name="Callout · Compliance Gateway")

# 7) G2 Agent Core (hero, 0.72x0.72 h=0.26)
iso_box(4.304, 2.695, 0.72, 0.72, 0.26, ACCENT, "Agent Core", hero=True)
callout(4.31, 3.57, "Agent Core", accent=True, anchor="tc", name="Callout · Agent Core")

# 8) G3 children — LLM 슬래브 3장 + SLM 2장 + Model Gateway
for i, t in enumerate([1.51, 1.39, 1.27]):
    iso_box(5.32, t, 0.36, 0.36, 0.11, MODEL, f"LLM slab {i+1}")
textbox(OX + 4.78, OY + 1.32, "LLM", 7, mix(TEXT, 0.75), w=0.5, h=0.15, name="LLM label")
for i, t in enumerate([1.84, 1.72]):
    iso_box(6.1, t, 0.27, 0.27, 0.11, MODEL, f"SLM slab {i+1}")
textbox(OX + 6.42, OY + 1.76, "SLM", 7, mix(TEXT, 0.75), w=0.5, h=0.15, name="SLM label")
iso_box(5.54, 1.66, 0.55, 0.55, 0.26, ACCENT, "Model Gateway")
# gateway 슬롯 (f-top::after: 14%~86% x, 40%~60% y 다크 바)
gx, gy = OX + 5.54, OY + 1.66
slot_pts = [(gx + iso(u, v)[0], gy + iso(u, v)[1])
            for u, v in [(0.077, 0.22), (0.473, 0.22), (0.473, 0.33), (0.077, 0.33)]]
freeform(slot_pts, fill=mix(ACCENT, 0.45, base=BG), line=None).name = "Model Gateway · slot"
leader(5.54, 1.66, 0.9, "Model Gateway")
callout(5.54, 0.76, "Model Gateway", anchor="bl", name="Callout · Model Gateway")

# 9) G4 tool cubes x3
iso_box(7.55, 1.63, 0.3, 0.3, 0.16, LAYER, "Tool cube 1")
iso_box(7.16, 1.85, 0.3, 0.3, 0.16, LAYER, "Tool cube 2")
iso_box(7.94, 1.85, 0.3, 0.3, 0.16, LAYER, "Tool cube 3")
leader(7.55, 1.63, 0.83, "Tools")
callout(7.55, 0.8, "Tools", anchor="bl", name="Callout · Tools")

# 10) G5 cylinders x4 + dict
iso_cyl(9.25, 2.59, 0.46, 0.22, STORE, "Store cyl 1")
iso_cyl(9.8, 2.79, 0.46, 0.22, STORE, "Store cyl 2")
iso_cyl(8.9, 2.83, 0.46, 0.22, STORE, "Store cyl 3")
iso_cyl(9.45, 3.04, 0.46, 0.22, STORE, "Store cyl 4")
iso_cyl(10.0, 3.19, 0.38, 0.16, DICT, "Dictionary cyl")
leader(9.48, 2.62, 0.6, "Stores")
callout(9.48, 2.0, "Stores", anchor="bl", name="Callout · Stores")

# 11) G6 Ingestion pipe — CAN 회전 근사
# CSS: left 5.5, top 4.6, 2.7x0.52, origin 0 50%, rotate(-26.57deg)
plen, pdia, ang = 2.7, 0.52, 26.57
pax, pay = OX + 5.5, OY + 4.6 + pdia / 2          # 회전 원점 (왼쪽 중앙)
ca, sa = math.cos(math.radians(-ang)), math.sin(math.radians(-ang))
pcx = pax + (plen / 2) * ca
pcy = pay + (plen / 2) * sa
pipe = shapes.add_shape(MSO_SHAPE.CAN, Inches(pcx - pdia / 2), Inches(pcy - plen / 2),
                        Inches(pdia), Inches(plen))
pipe.rotation = 90 - ang    # 세로 CAN 축을 -26.57° 상행(우상향) 방향으로
pipe.fill.solid(); pipe.fill.fore_color.rgb = mix(DUCT, 0.6)
pipe.line.color.rgb = mix(DUCT, 0.6, base=(0x1F, 0x1F, 0x1F)); pipe.line.width = Pt(1.0)
pipe.shadow.inherit = False; pipe.name = "Ingestion pipe"
# 파이프 축 대시 흐름선
fx1, fy1 = pax + 0.5 * ca, pay + 0.5 * sa
fx2, fy2 = pax + 2.45 * ca, pay + 2.45 * sa
fl = shapes.add_connector(1, Inches(fx1), Inches(fy1), Inches(fx2), Inches(fy2))
fl.line.color.rgb = mix((0xFF, 0xFF, 0xFF), 0.5, base=DUCT); fl.line.width = Pt(1.6)
ln = fl.line._get_or_add_ln()
ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'sysDash'}))
fl.shadow.inherit = False; fl.name = "Ingestion pipe · flow"
callout(7.62, 4.14, "Ingestion Pipeline", anchor="bl", name="Callout · Ingestion Pipeline")

# 12) 푸터
textbox(0.7, 7.08, "KORI Answers · XENOIMPACT", 8, mix(TEXT, 0.5), bold=False,
        w=4, h=0.2, name="Footer")

out = os.path.join(os.path.dirname(__file__), "..", "output",
                   "20260822_kori_architecture_hero_editable_v1.pptx")
out = os.path.abspath(out)
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved:", out)
